"""
extract_text.py — Extract Japanese text from course materials (PDF/pptx).

Strategy:
- Only include core course material types (授業シート, 補足資料, デモ, ガイダンス, 資料 etc.)
- Exclude administrative files (カルテ, コマシラバス, シラバス, 担当講師...紹介, etc.)
- Redact PII (names, emails, phones, student IDs, URLs)
- Output one .txt per course folder into scripts/extraction/raw_text/
"""

from __future__ import annotations

import os
import re
import sys
from pathlib import Path
from typing import Iterable

import pdfplumber
from pptx import Presentation

SOURCE_ROOT = Path("/Users/shanlei/Desktop/学校/AAA専門学校東京テクニカルカレッジ")
OUTPUT_ROOT = Path("/Users/shanlei/Desktop/ds-tango/scripts/extraction/raw_text")

# Filename keywords -- only keep files whose name contains one of these
INCLUDE_KEYWORDS = [
    "授業シート",
    "授業スライド",
    "授業資料",
    "スライド",
    "補足資料",
    "ガイダンス",
    "デモ",
    "解説",      # 課題解説 etc, but careful: カルテ解説 also matches -> use EXCLUDE first
    "資料",
    "教材",
    "テキスト",
    "ハンドアウト",
    "演習",
    "実習",
    "講義",
]

# These take precedence — if a filename contains any, skip it
EXCLUDE_KEYWORDS = [
    "カルテ",          # quizzes / homework worksheets
    "コマシラバス",     # syllabus admin
    "シラバス",
    "紹介",            # teacher/student intros
    "担当",
    "出席",
    "成績",
    "名簿",
    "アンケート",
    "課題提出",
    "解答",            # answer keys (may contain student-specific info)
]

# Extensions per PRD
INCLUDE_EXTS = {".pdf", ".pptx", ".ppt"}

# ---------------- Redaction patterns ----------------
EMAIL_RE = re.compile(r"\b[\w.+-]+@[\w.-]+\.[A-Za-z]{2,}\b")
URL_RE = re.compile(r"https?://\S+|www\.\S+")
PHONE_RE = re.compile(r"\b0\d{1,4}[-(]?\d{1,4}[-)]?\d{3,4}\b")
# Student-ID-like: 6-12 digit run (course materials usually don't have such bare numbers).
# We are conservative: only redact when surrounded by ID context, OR very long bare digit runs.
LONG_DIGITS_RE = re.compile(r"\b\d{7,}\b")
# Honorific name patterns
HONORIFIC_RE = re.compile(r"[一-龥々ヶー][一-龥々ヶーぁ-んァ-ヴ]{0,5}(?:先生|さん|くん|君|ちゃん|教員|教官|講師)")
# Date-stamp / page-number style stuff we don't want as text
PAGE_FOOTER_RE = re.compile(r"^\s*(?:Page|ページ|p\.?)\s*\d+\s*(?:/\s*\d+)?\s*$", re.IGNORECASE)
COPYRIGHT_RE = re.compile(r"(?:©|Copyright|\(c\))\s*[\w぀-ヿ一-鿿 .,'-]{0,80}", re.IGNORECASE)
# Common校名 (school names) that we want stripped from headers/footers
SCHOOL_NAMES = [
    "東京テクニカルカレッジ",
    "AAA専門学校",
    "TTC",
]

# Lines containing only these keywords (very short admin headers) are dropped
ADMIN_HEADER_HINTS = [
    "学籍番号",
    "氏名",
    "受講者",
    "提出日",
    "出席",
    "コマシラバス",
]


def should_include(filename: str) -> bool:
    """Return True if filename matches an include keyword and not an exclude keyword."""
    name = filename
    for kw in EXCLUDE_KEYWORDS:
        if kw in name:
            return False
    for kw in INCLUDE_KEYWORDS:
        if kw in name:
            return True
    return False


def redact_line(line: str) -> str | None:
    """Apply redaction. Return cleaned line, or None to drop entirely."""
    if not line.strip():
        return None
    # Strip page footers and short admin headers
    if PAGE_FOOTER_RE.match(line):
        return None
    if any(h in line and len(line.strip()) < 30 for h in ADMIN_HEADER_HINTS):
        return None
    # Drop copyright lines
    if COPYRIGHT_RE.search(line):
        line = COPYRIGHT_RE.sub("", line)
    # Email / URL / phone -> remove
    line = EMAIL_RE.sub("", line)
    line = URL_RE.sub("", line)
    line = PHONE_RE.sub("", line)
    # Long bare digit runs -> remove
    line = LONG_DIGITS_RE.sub("", line)
    # Names with honorifics -> remove the whole token
    line = HONORIFIC_RE.sub("", line)
    # School names
    for sn in SCHOOL_NAMES:
        line = line.replace(sn, "")
    line = line.strip()
    if len(line) < 2:
        return None
    return line


def extract_pdf(path: Path) -> Iterable[str]:
    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                for raw in text.splitlines():
                    cleaned = redact_line(raw)
                    if cleaned:
                        yield cleaned
    except Exception as e:
        print(f"  [warn] PDF failed: {path.name} ({e})", file=sys.stderr)


def extract_pptx(path: Path) -> Iterable[str]:
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if not text:
                        continue
                    for raw in text.splitlines():
                        cleaned = redact_line(raw)
                        if cleaned:
                            yield cleaned
    except Exception as e:
        print(f"  [warn] PPTX failed: {path.name} ({e})", file=sys.stderr)


def process_course(course_dir: Path) -> tuple[int, int]:
    """Process one course folder. Returns (files_used, lines_written)."""
    out_path = OUTPUT_ROOT / f"{course_dir.name}.txt"
    files_used = 0
    lines_written = 0
    with out_path.open("w", encoding="utf-8") as out:
        out.write(f"# Source: {course_dir.name}\n\n")
        for p in sorted(course_dir.rglob("*")):
            if not p.is_file():
                continue
            if p.suffix.lower() not in INCLUDE_EXTS:
                continue
            if not should_include(p.name):
                continue
            # Always extract — but never write filename (to avoid leaking instructor info that may live in filename)
            files_used += 1
            out.write(f"\n## FILE_{files_used:03d}\n")
            extractor = extract_pdf if p.suffix.lower() == ".pdf" else extract_pptx
            for line in extractor(p):
                out.write(line + "\n")
                lines_written += 1
    return files_used, lines_written


def main():
    OUTPUT_ROOT.mkdir(parents=True, exist_ok=True)
    course_dirs = sorted([d for d in SOURCE_ROOT.iterdir() if d.is_dir() and re.match(r"^\d\d_", d.name)])
    print(f"Found {len(course_dirs)} course folders.")
    grand_files = 0
    grand_lines = 0
    for d in course_dirs:
        files_used, lines = process_course(d)
        grand_files += files_used
        grand_lines += lines
        print(f"  {d.name}: {files_used} files, {lines} lines")
    print(f"\nTotal: {grand_files} files, {grand_lines} lines")
    print(f"Output: {OUTPUT_ROOT}")


if __name__ == "__main__":
    main()
